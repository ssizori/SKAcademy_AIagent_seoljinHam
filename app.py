import streamlit as st
import pandas as pd
import numpy as np
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO

# 1. 웹 페이지 기본 설정
st.set_page_config(
    page_title="SK hynix AI 8D Report Platform",
    page_icon="🔬",
    layout="wide"
)

# 2. 헤더 및 타이틀
st.title("🔬 제조 및 반도체 공정 최적화 AI 에이전트 (V4.4)")
st.markdown("---")

# 3. 사이드바 - 제어 패널
st.sidebar.header("🎛️ 공정 제어 패널")

process_type = st.sidebar.selectbox(
    "1. 모니터링할 핵심 공정을 선택하세요",
    ["Photolithography", "Etching", "Deposition", "Cleaning", "Diffusion"]
)

# 공정별 메타데이터 세팅
process_meta = {
    "Photolithography": {
        "defects": ["Overlay Misalignment", "Critical Dimension (CD) Error", "PR Peeling"],
        "metrics": ["노광 에너지", "포커스 Margin", "오버레이 에러"],
        "targets": [150.5, 0.25, 14.0],
        "units": ["mJ", "µm", "nm"],
        "ucl": 18.0, "lcl": 10.0,
        "recipe": "ArF Immersion Baseline / Target CD: 45nm"
    },
    "Etching": {
        "defects": ["Under Etch", "Over Etch / Profile Bowing", "Selectivity Loss"],
        "metrics": ["RF 파워", "가스 압력", "식각 CD"],
        "targets": [850.0, 15.2, 38.0],
        "units": ["W", "mTorr", "nm"],
        "ucl": 42.0, "lcl": 34.0,
        "recipe": "High Aspect Ratio Oxide Etch / Target Depth: 2.5µm"
    },
    "Deposition": {
        "defects": ["Step Coverage Failure", "Thickness Non-uniformity", "Particle Contamination"],
        "metrics": ["증착 온도", "챔버 압력", "박막 두께 편차"],
        "targets": [450.0, 2.5, 15.0],
        "units": ["°C", "Torr", "Å"],
        "ucl": 25.0, "lcl": 5.0,
        "recipe": "Spectral Atomic Layer Deposition (ALD) High-k Baseline"
    },
    "Cleaning": {
        "defects": ["Water Mark Residual", "Under Cleaning / Native Oxide", "Pattern Collapse"],
        "metrics": ["매엽 가동 시간", "화학액 농도", "잔류 파티클 수"],
        "targets": [35.0, 5.5, 10.0],
        "units": ["sec", "%", "ea"],
        "ucl": 30.0, "lcl": 0.0,
        "recipe": "DHF / SC1 Single Wafer Cleaning Recipe"
    },
    "Diffusion": {
        "defects": ["Junction Depth Shift", "Sheet Resistance (Rs) Variation", "Slip Line Defect"],
        "metrics": ["메인 튜브 온도", "공정 압력", "면저항 변동치"],
        "targets": [980.0, 760.0, 50.0],
        "units": ["°C", "Torr", "ohm/sq"],
        "ucl": 55.0, "lcl": 45.0,
        "recipe": "Well Drive-in High Temp Oxidation / Target Rs: 50 ohm/sq"
    }
}

current_meta = process_meta[process_type]
defect_type = st.sidebar.selectbox("2. 감지된 불량 유형을 선택하세요", current_meta["defects"])
worker_note = st.sidebar.text_area("3. 작업자 특이사항 및 설비 로그 노트", placeholder="예: 챔버 헌팅 현상 관측 등")

st.sidebar.markdown("---")
run_agent = st.sidebar.button("AI 에이전트 분석 & 8D 생성", type="primary")

# 4. 메인 화면 - 탭 구성
tab1, tab2 = st.tabs(["📊 실시간 센서 트렌드", "🤖 SK하이닉스 양식 8D Report & 이메일"])

with tab1:
    st.subheader(f"⚡ {process_type} 공정 실시간 데이터 모니터링")
    st.info(f"📋 현재 라인 레시피 가이드라인: {current_meta['recipe']}")
    
    target_metric_name = current_meta["metrics"][2]
    target_val = current_meta["targets"][2]
    unit_val = current_meta["units"][2]
    
    np.random.seed(42)
    time_points = 25
    base_signal = np.random.randn(time_points) * 0.8 + target_val
    anomaly_start = 15
    
    if defect_type in ["Overlay Misalignment", "Over Etch / Profile Bowing", "Thickness Non-uniformity", "Junction Depth Shift"]:
        base_signal[anomaly_start:] += np.linspace(0, 6, time_points - anomaly_start)
    else:
        base_signal[anomaly_start:] -= np.linspace(0, 5, time_points - anomaly_start)
        
    latest_val = base_signal[-1]
    is_out = latest_val > current_meta["ucl"] or latest_val < current_meta["lcl"]
    
    col1, col2, col3 = st.columns(3)
    col1.metric(
        label=f"{current_meta['metrics'][0]} ({current_meta['units'][0]})", 
        value=f"{current_meta['targets'][0] + np.random.uniform(-0.5, 0.5):.1f} {current_meta['units'][0]}", 
        delta="정상 수렴"
    )
    col2.metric(
        label=f"{current_meta['metrics'][1]} ({current_meta['units'][1]})", 
        value=f"{current_meta['targets'][1] + np.random.uniform(-0.01, 0.01):.3f} {current_meta['units'][1]}", 
        delta="정상 마진 내"
    )
    
    if is_out:
        col3.metric(
            label=f"🚨 현재 {target_metric_name} (목표: {target_val}{unit_val})", 
            value=f"{latest_val:.1f} {unit_val}", 
            delta="⚠️ 관리 한계선 이탈!", 
            delta_color="inverse"
        )
    else:
        col3.metric(
            label=f"🟢 현재 {target_metric_name} (목표: {target_val}{unit_val})", 
            value=f"{latest_val:.1f} {unit_val}", 
            delta="정상 내부 구동"
        )

    df_chart = pd.DataFrame({
        target_metric_name: base_signal,
        "관리 상한선 (UCL)": [current_meta["ucl"]] * time_points,
        "관리 하한선 (LCL)": [current_meta["lcl"]] * time_points,
        "공정 목표 수치 (Target)": [target_val] * time_points
    })
    st.line_chart(df_chart, color=["#FF4B4B", "#FFAA00", "#FFAA00", "#00A0FF"])

with tab2:
    if run_agent:
        note_content = worker_note if worker_note != "" else "공정 테이블 내 일시적 변동성 관측"
        doc_no = f"2026-SKH-{process_type[:3].upper()}-007"
        current_date = time.strftime('%Y년 %m월 %d일')
        
        # 📌 1. 웹 대책서 순서 정렬 및 단일 구조화 (D1 ~ D8 순차 진행)
        st.markdown(f"<h1 style='text-align: center; color: #000000; font-weight: bold;'>SK하이닉스 {process_type} [{defect_type}] 8D Report</h1>", unsafe_allow_html=True)
        st.write("")
        
        # 메타 정보 테이블
        meta_df = pd.DataFrame([
            {"정보 항목": "D1. 담당 팀 (Team)", "세부 내용": f"양산기술 {process_type}팀 (담당 엔지니어: 함설진 TL)", "정보 항목 ": "문서번호 / 통보일자", "세부 내용 ": f"{doc_no} / {current_date}"},
            {"정보 항목": "대상 부품/공정", "세부 내용": f"{process_type} Main Line", "정보 항목 ": "발견된 불량 유형", "세부 내용 ": defect_type},
            {"정보 항목": "실시간 측정 파라미터", "세부 내용": f"{target_metric_name} (측정치: {latest_val:.1f} {unit_val} / 목표: {target_val} {unit_val})", "정보 항목 ": "수신 참조 부서", "세부 내용 ": "본사 품질보증팀"}
        ])
        st.table(meta_df.set_index("정보 항목"))
        
        # 순차적 단계 통합 서술
        with st.container():
            st.markdown("### 🟠 D2. 불량 현상 정의 (Problem Description)")
            st.markdown(f"""
            * {process_type} 공정 모니터링 중 {target_metric_name}의 관리 한계값 이탈 현상 실시간 포착.
            * **현장 특이사항 및 설비 로그:** {note_content}
            """)
            st.markdown("---")
            
            st.markdown("### 🟠 D3. 임시 봉쇄 조치 (Containment Action)")
            st.markdown("""
            * SPC 연동 설비 실시간 고유 인터락 가동 및 챔버 구동 즉시 Hold 조치.
            * 전후 인접 인라인 운송 Lot 전수 검사 지정 및 격리 처리를 통한 불량 유출 차단.
            """)
            st.markdown("---")
            
            st.markdown("### 🟠 D4. 근본 원인 분석 (Root Cause Analysis)")
            st.markdown("""
            * 센서 데이터 추적 결과 15번 시퀀스 시점부터 급격한 경향성 이탈(Drift) 트렌드가 기록됨.
            * 하드웨어 부품 피로도 누적 및 물리적 마모에 따른 파라미터 헌팅 현상이 지배적 원인으로 분석 완료.
            """)
            st.markdown("---")
            
            st.markdown("### 🟠 D5~D6. 영구 시정 조치 수립 및 유효성 검증 (Corrective Actions)")
            st.markdown("""
            * APC 피드백 제어 매칭 알고리즘 업데이트 및 보정 계수 수정 단행.
            * 정밀 하드웨어 교정을 통한 물리 레시피 기준 마진 원복 및 표준화 세팅 단행.
            """)
            st.markdown("---")
            
            st.markdown("### 🟠 D7. 재발 방지 대책 수립 (Preventive Action)")
            st.markdown("""
            * OCAP 품질 관리 차트 이상 징후 자동 감시 인터락 한계선을 기존 대비 15% 타이트하게 변경 적용.
            * 이상 발생 설비 대상 파라미터 감시 모니터링 주기를 단축 설정하여 예방 보전 체계 고도화.
            """)
            st.markdown("---")
            
            st.markdown("### 🟠 D8. 개선 활동 완료 검토 및 팀 승인 (Congratulate Your Team)")
            st.markdown(f"""
            * **최종 확인 결과:** ☑ Closed &nbsp;&nbsp; ☐ Unclosed
            * 제반 센서 추적 수치가 안정권 내로 완벽하게 복귀하여 가동 지속 정상 판정 확인.
            
            <p style='text-align: right; font-weight: bold;'>승인자: 양산기술 {process_type} 담당 함설진 TL (인) &nbsp;&nbsp;&nbsp;&nbsp; 승인일자: {current_date}</p>
            """, unsafe_allow_html=True)


        # 📌 2. PPTX 서식 완전 수정 (D2 단락 쪼개기 기법 적용으로 글씨 크기 12pt 완벽 고정)
        def generate_sk_pptx():
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # 테마 상단 SK 오렌지 라인
            top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = RGBColor(240, 90, 40)
            top_bar.line.fill.background()
            
            # 📌 수정 조치: 기존 주황색 조그만 문구를 삭제/통합하고 상단에 검은색 큰 텍스트로 타이틀 일원화
            header_text_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11.0), Inches(0.6))
            header_p = header_text_box.text_frame.paragraphs[0]
            header_p.text = f"SK하이닉스 {process_type} [{defect_type}] 8D Report"
            header_p.font.size = Pt(24)
            header_p.font.bold = True
            header_p.font.color.rgb = RGBColor(0, 0, 0) # 검은색 변경
            
            # 테마 하단 SK 오렌지 라인
            bottom_bar = slide.shapes.add_shape(1, Inches(0), Inches(7.2), Inches(13.333), Inches(0.04))
            bottom_bar.fill.solid()
            bottom_bar.fill.fore_color.rgb = RGBColor(240, 90, 40)
            bottom_bar.line.fill.background()
            
            # 📌 수정 조치: 하단 문구를 간결하게 "8D Report"로만 축소 변경
            footer_box = slide.shapes.add_textbox(Inches(5.0), Inches(7.23), Inches(3.33), Inches(0.25))
            footer_p = footer_box.text_frame.paragraphs[0]
            footer_p.text = "8D Report"
            footer_p.font.size = Pt(10)
            footer_p.font.bold = True
            footer_p.font.color.rgb = RGBColor(128, 128, 128)
            footer_p.alignment = PP_ALIGN.CENTER
            
            # SK hynix 우측 상단 텍스트 로고 영역 배치
            logo_box = slide.shapes.add_textbox(Inches(11.0), Inches(0.4), Inches(1.8), Inches(0.5))
            logo_p = logo_box.text_frame.paragraphs[0]
            logo_p.text = "SK hynix"
            logo_p.font.size = Pt(18)
            logo_p.font.bold = True
            logo_p.font.color.rgb = RGBColor(240, 90, 40)
            logo_p.alignment = PP_ALIGN.RIGHT
            
            # 데이터 표 초기화 수립 (D1~D8 순차 배치 구조화)
            rows, cols = 8, 2
            left, top, width, height = Inches(0.6), Inches(1.3), Inches(12.13), Inches(5.6)
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            table.columns[0].width = Inches(2.8)
            table.columns[1].width = Inches(9.33)
            
            # 각 행별 바인딩 데이터 세팅
            headers = [
                ("D1. 담당 팀 및 메타정보", [f"발행팀: 양산기술 {process_type}팀 / 담당 엔지니어: 함설진 TL", f"문서번호: {doc_no} / 발행일자: {current_date}"]),
                
                # 📌 문제의 D2 타겟 구간: 리스트 형태로 문장을 쪼개어 각각 단락 객체로 제어
                ("D2. 불량 현상 정의", [
                    f"공정명 및 불량유형: {process_type} Main Line / {defect_type}",
                    f"실시간 센서 측정치: {latest_val:.1f} {unit_val} (Target 기준수치: {target_val} {unit_val})",
                    f"엔지니어 현장로그: {note_content}"
                ]),
                
                ("D3. 임시 봉쇄 조치", ["SPC 인터락 연동을 통한 설비 즉시 셧다운(Hold) 조치 구동 완료", "전후 인접 인라인 설비 대상 추적 Lot 전량 격리 지정 완료"]),
                ("D4. 근본 원인 분석", ["15번 시퀀스 시점 센서 데이터 분석 결과 경향성 이탈(Drift) 현상 포착", "하드웨어 장비 내부 부품의 피로도 누적 및 마모에 따른 변동으로 판명"]),
                ("D5~D6. 영구 시정 조치", ["APC 피드백 제어 보정 계수 모델 업데이트 반영 완료", "하드웨어 물리 레시피 기준 마진 정밀 캘리브레이션 셋업 원복"]),
                ("D7. 재발 방지 대책", ["OCAP 품질 관리 차트 이상 감시 알람 한계선 15% 타이트화 설정 적용", "설비 모니터링 샘플링 주기를 단축 설정하여 상시 감시 체계 고도화"]),
                ("D8. 개선 활동 완료 검토", ["최종 검토 결과: Closed (수치 정상 범위 안착 가동 확인 완료)", f"최종 승인 엔지니어: 양산기술팀 함설진 TL (서명 날인 생략)"])
            ]
            
            for idx, (label, val_list) in enumerate(headers):
                # 좌측 설명 열 세팅 (오렌지 톤 연출)
                cell_lbl = table.cell(idx, 0)
                cell_lbl.text = label
                cell_lbl.fill.solid()
                cell_lbl.fill.fore_color.rgb = RGBColor(255, 238, 230)
                p_lbl = cell_lbl.text_frame.paragraphs[0]
                p_lbl.font.bold = True
                p_lbl.font.size = Pt(12)
                p_lbl.font.color.rgb = RGBColor(210, 70, 20)
                
                # 우측 데이터 내용 열 세팅 (깔끔한 연회색 바탕)
                cell_val = table.cell(idx, 1)
                cell_val.fill.solid()
                cell_val.fill.fore_color.rgb = RGBColor(248, 249, 250)
                
                # 📌 핵심 해결책: 기존 한꺼번에 집어넣던 형식을 버리고 개별 단락을 명시적으로 생성하여 각각 12pt 지정
                tf_val = cell_val.text_frame
                tf_val.clear() # 기본 생성 단락 초기화
                
                for line_idx, line_text in enumerate(val_list):
                    if line_idx == 0:
                        p_line = tf_val.paragraphs[0]
                    else:
                        p_line = tf_val.add_paragraph()
                    
                    p_line.text = f"• {line_text}"
                    p_line.font.size = Pt(12) # 예외 없이 강제 고정 픽스 적용
                    p_line.font.color.rgb = RGBColor(0, 0, 0)
                
            binary_output = BytesIO()
            prs.save(binary_output)
            return binary_output.getvalue()

        pptx_data = generate_sk_pptx()
        st.write("")
        
        # 📌 3. 다운로드 버튼 이름 동적 변경 유지
        st.download_button(
            label=f"📊 {process_type} 공정 [{defect_type}] 8D Report 다운로드",
            data=pptx_data,
            file_name=f"SKH_8D_Report_{process_type}_{defect_type.replace(' ', '_')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        # 5. 이메일 자동 작성 섹션
        st.markdown("---")
        st.markdown("### ✉️ 정형 양식 기반 이메일 자동 완성")
        
        email_tab1, email_tab2 = st.tabs(["👨‍✈️ 선임 엔지니어 보고용 이메일", "🏢 고객사 송부용 이메일"])
        
        with email_tab1:
            st.code(f"""
제목: [공정보고] {process_type} 공정 {defect_type} 불량 발생에 따른 8D 조치 현황 보고의 건

양산기술 {process_type} 담당 함설진 TL입니다.

금일 메인 라인 모니터링 중 발생한 {process_type} 공정 내 {defect_type} 불량에 대한 원인 분석 및 조치 사항을 아래와 같이 요약 보고 드립니다.

1. 발생 공정 및 불량명: {process_type} / {defect_type}
2. 주요 센서 이상 징후: {target_metric_name} 실시간 수치 {latest_val:.1f} {unit_val} (목표 수치: {target_val} {unit_val} 대비 한계치 이탈)
3. 현장 조치 내용 (D3~D6):
   - 해당 설비 즉시 홀드 및 인터락 마진 한계 세팅 조정 완료
   - 엔지니어링 분석을 통한 레시피 표준화 및 APC 파라미터 원복 완료
4. 특이사항: {note_content}

자세한 수율 영향성 및 D8 개선 활동 완료 여부 검토를 포함한 8D Report 서명본은 첨부 문서를 확인해 주시기 바랍니다.

감사합니다.
양산기술 {process_type}팀 함설진 TL 드림.
            """, language="text")
                
        with email_tab2:
            st.code(f"""
제목: [Technical Notification] {process_type} 공정 이상 제어 및 안정화 완료에 따른 안내의 건

안녕하십니까, 양산기술 {process_type} 담당 함설진 TL입니다.

언제나 저희 제품을 신뢰해 주시고 협력해 주시는 고객사의 무궁한 발전을 기원합니다.

최근 {process_type} 공정 라인에서 일시적인 파라미터 드리프트로 인해 [{defect_type}] 관련 미세 마진 저하 우려가 감지되었으나, 당사의 실시간 AI 에이전트 시스템 및 표준품질 양식 트러블 슈팅 절차에 의거하여 즉각적인 시정 조치(Corrective Action)를 완료하였음을 안내해 드립니다.

[공정 안정화 요약]
- 대상 공정: {process_type}
- 조치 사항: 센서 제어 조건 표준화 정밀 보정 및 예방 보전(PM) 강화 적용
- 수율 및 품질 영향성: 임시 봉쇄 및 내부 검증 결과, 출하 제품 품질 마진에는 영향 없음이 최종 확인 및 8D Report 클로징(Closed)되었습니다.

앞으로도 철저한 공정 표준화 조치를 통해 안정적인 고품질 반도체 공급을 약속드립니다. 본 건과 관련하여 추가 문의 사항이 있으신 경우 언제든 편하게 연락해 주시기 바랍니다.

감사합니다.
양산기술 {process_type}팀 함설진 TL 드림.
            """, language="text")
            
    else:
        st.info("👈 왼쪽 사이드바에서 [AI 에이전트 분석 & 8D 생성] 버튼을 누르면 8D Report가 빌드됩니다.")

# 5. 하단 푸터
st.markdown("---")
st.caption("© 2026 SK hynix Smart Manufacturing AI. 본 시스템은 사내 보안 가이드라인을 준수합니다.")
